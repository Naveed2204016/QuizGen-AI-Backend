import io
import re
from dataclasses import dataclass

import pymupdf
import pytesseract
from pptx import Presentation
from pypdf import PdfReader

from app.core.config import get_settings


@dataclass
class DocumentSection:
    number: int
    label: str
    text: str


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    return re.sub(r"[ \t]+", " ", re.sub(r"\n{3,}", "\n\n", text)).strip()


def _has_meaningful_text(text: str) -> bool:
    """Reject scanner watermarks and similarly content-free extracted pages."""
    without_scanner_labels = re.sub(
        r"(?im)^\s*(?:scanned\s+by\s+camscanner|camscanner)\s*$", "", text
    ).strip()
    # Count letters from any Unicode script so Bangla and other languages are valid.
    words = re.findall(r"[^\W\d_]{2,}", without_scanner_labels, flags=re.UNICODE)
    return (
        len(without_scanner_labels) >= 20
        and len(words) >= 5
        and len(set(map(str.lower, words))) >= 4
    )


def _ocr_pdf_page(document: pymupdf.Document, page_index: int) -> str:
    """Render one PDF page and extract its text with local Tesseract OCR."""
    settings = get_settings()
    if settings.tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

    page = document.load_page(page_index)
    pixmap = page.get_pixmap(
        dpi=settings.ocr_dpi,
        colorspace=pymupdf.csRGB,
        alpha=False,
    )
    image = pixmap.pil_image()
    try:
        return _clean(pytesseract.image_to_string(image, lang=settings.ocr_language))
    finally:
        image.close()


def _extract_pdf(content: bytes) -> list[DocumentSection]:
    reader = PdfReader(io.BytesIO(content))
    sections = [
        DocumentSection(index, f"Page {index}", _clean(page.extract_text() or ""))
        for index, page in enumerate(reader.pages, 1)
    ]

    settings = get_settings()
    unreadable_indexes = [
        index for index, section in enumerate(sections) if not _has_meaningful_text(section.text)
    ]
    if not settings.ocr_enabled or not unreadable_indexes:
        return sections

    try:
        with pymupdf.open(stream=content, filetype="pdf") as document:
            for page_index in unreadable_indexes:
                sections[page_index].text = _ocr_pdf_page(document, page_index)
    except pytesseract.TesseractNotFoundError as exc:
        raise ValueError(
            "This PDF needs OCR, but Tesseract is not installed or configured."
        ) from exc
    except pytesseract.TesseractError as exc:
        raise ValueError(f"OCR failed: {exc}") from exc

    return sections


def extract_document(content: bytes, extension: str) -> list[DocumentSection]:
    if extension == ".pdf":
        sections = _extract_pdf(content)
    elif extension == ".pptx":
        presentation = Presentation(io.BytesIO(content))
        sections = []
        for index, slide in enumerate(presentation.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(" | ".join(cell.text for cell in row.cells))
            sections.append(DocumentSection(index, f"Slide {index}", _clean("\n".join(parts))))
    else:
        raise ValueError("Only PDF and PPTX files are supported")

    useful = [section for section in sections if _has_meaningful_text(section.text)]
    if not useful:
        raise ValueError(
            "No readable study text was found, even after OCR. Try a clearer scan or check the OCR language."
        )
    return useful
